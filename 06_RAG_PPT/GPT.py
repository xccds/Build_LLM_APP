from langchain_openai import ChatOpenAI,OpenAIEmbeddings
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores.utils import filter_complex_metadata
from langchain_community.vectorstores import Chroma
from langchain.tools.retriever import create_retriever_tool
from langchain import hub
from langchain.agents import AgentExecutor, create_openai_functions_agent
from langchain_openai import ChatOpenAI
from langchain_community.tools.tavily_search import TavilySearchResults

from openai import OpenAI
from io import BytesIO
from pptx import Presentation
import json
import requests

from dotenv import load_dotenv
load_dotenv()


class PowerpointAI:

    def __init__(self):
        self.client = OpenAI()
        self.model = ChatOpenAI(temperature=0.1, model_name="gpt-3.5-turbo")
        self.text_splitter = RecursiveCharacterTextSplitter(chunk_size=1024, chunk_overlap=100)
        self.embeddings = OpenAIEmbeddings()
        self.search_tool = TavilySearchResults()
        instructions = """You are a power point presentation specialist."""
        base_prompt = hub.pull("langchain-ai/openai-functions-template")
        self.prompt = base_prompt.partial(instructions=instructions)
        self.tools = [self.search_tool]

    def ingest(self, file_path: str):
        docs = PyPDFLoader(file_path=file_path).load()
        chunks = self.text_splitter.split_documents(docs)
        chunks = filter_complex_metadata(chunks)
        vector_store = Chroma.from_documents(documents=chunks, embedding=self.embeddings)
        self.retriever = vector_store.as_retriever(search_type = "mmr", search_kwargs={
            "k": 6,
            "fetch_k": 20,
            "include_metadata": True
            },
        )
        self.retriever_tool = create_retriever_tool(
                self.retriever,
                "uploaded_files_search",
                "Search for information about uploaded files, for any question about uploaded files you should use this tool",
            )
        self.tools.append(self.retriever_tool)
        self.agent = create_openai_functions_agent(self.model, self.tools, self.prompt)



    def ask(self, questioin: str,num_slides: int, search:bool=False, rag:bool=False):

        base_prompt = f"question: {questioin}"
        if search:
            base_prompt += " using web search tool"
        if rag:
            base_prompt += " using uploaded file"
        instructions = f"""
        you should create a PowerPoint presentation to answer the above question. Make it {num_slides}  slides.
        Each slide should have a title and content.
        step 1: think and write an outline for each title.
        step 2: with the content, expand on each of the title of your outline using at least 3 to 5 sentences. 
        You can consider elaborating on the key ideas, offering supporting examples and explaining any details that you think would enhance the audience's understanding on the topic.
        step 3: Structure the output information as a JSON as follows.
        Your answer should only contain the JSON - no markdown formatting.
        """
        format_prompt ="""
        Example:
        {"slides": [
        {"title": "title for slide 1", "content": "Content for slide 1"},
        {"title": "title for slide 2", "content": "Content for slide 2"},
        {"title": "title for slide 3", "content": "Content for slide 3"},
        ]}
        """
        final_prompt = base_prompt + instructions + format_prompt

        agent_executor = AgentExecutor(
            agent=self.agent,
            tools=self.tools,
            verbose=False)
        result = agent_executor.invoke({"input": final_prompt})
        self.question = questioin
        self.json_output = result["output"]

    def json2dict(self):
        self.dict_output = json.loads(self.json_output)['slides']

    def image_generator(self,prompt):
        try:
            # Create an image by using the image generation API
            generation_response = self.client.images.generate(
                model="dall-e-3",
                prompt=prompt,
                size='1024x1024',
                n=1
            )
            image_url = generation_response.data[0].url # extract image URL from response
            image = requests.get(image_url).content  # download the image
            generated_image = BytesIO(image)
            return generated_image
        
        except self.client.error.InvalidRequestError as err:
            print(err)

    def create_ppt(self,topic):
        titles = [slide['title'] for slide in self.dict_output]
        contents = [slide['content'] for slide in self.dict_output]
        prs = Presentation()
        title_side_layout = prs.slide_layouts[0]
        two_col_slide_layout = prs.slide_layouts[8]

        title_slide = prs.slides.add_slide(title_side_layout)
        title_slide.shapes.title.text = topic

        for slide_title, slide_content in zip(titles, contents):
            slide = prs.slides.add_slide(two_col_slide_layout)

            # add title
            slide.shapes.title.text = slide_title

            # add image 
            image = self.image_generator(slide_content)
            slide.shapes.placeholders[1].insert_picture(image)

            # Add text 
            slide.shapes.placeholders[2].text = slide_content
        
        prs.save("presentation.pptx")
        return "presentation.pptx"

# if __name__ == "__main__":
#     powerpoint_ai = PowerpointAI()
#     powerpoint_ai.ingest("test.pdf")
#     powerpoint_ai.ask("What is critical thinking based on uploaded files?",num_slides=4)
#     powerpoint_ai.json2dict()
#     powerpoint_ai.create_ppt("What is critical thinking")
